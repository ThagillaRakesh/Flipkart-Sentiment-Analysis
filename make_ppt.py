"""
make_ppt.py — Generate an impressive project review PPT for Opinion Miner
Run: python3 make_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # slide background
DARK_CARD  = RGBColor(0x16, 0x2A, 0x3E)   # card background
TEAL       = RGBColor(0x00, 0xD4, 0xAA)   # accent / highlight
CORAL      = RGBColor(0xFF, 0x6B, 0x6B)   # warning / negative
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)   # gold star
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xC5, 0xD5, 0xE8)
MID_BLUE   = RGBColor(0x1E, 0x3A, 0x5F)
ORANGE     = RGBColor(0xFF, 0xA5, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ───────────────────────────────────────────────────────────

def bg(slide, colour=NAVY):
    """Fill slide background with a solid colour."""
    bg_ = slide.background
    fill = bg_.fill
    fill.solid()
    fill.fore_color.rgb = colour


def rect(slide, x, y, w, h, fill_colour, alpha=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, colour=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return txBox


def add_para(tf, text, size=16, bold=False, colour=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=6):
    """Append a paragraph to an existing text frame."""
    from pptx.util import Pt as pPt
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = pPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour


def label_box(slide, text, x, y, w, h, bg_col, text_col=WHITE, size=14, bold=True):
    """Filled rounded box with centred text (simulated with plain rect)."""
    rect(slide, x, y, w, h, bg_col)
    txt(slide, text, x, y, w, h, size=size, bold=bold, colour=text_col,
        align=PP_ALIGN.CENTER)


def divider(slide, y, colour=TEAL, thickness=0.04):
    """Horizontal divider line."""
    rect(slide, 0, y, 13.33, thickness, colour)


def heading(slide, title, subtitle=None):
    """Standard slide heading strip."""
    rect(slide, 0, 0, 13.33, 1.1, DARK_CARD)
    divider(slide, 1.06, TEAL, 0.06)
    txt(slide, title, 0.3, 0.1, 12.7, 0.7,
        size=30, bold=True, colour=TEAL, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.72, 12.7, 0.38,
            size=14, colour=LIGHT_GREY, align=PP_ALIGN.LEFT)


def bullet_card(slide, items, x, y, w, h, title=None,
                bg_col=DARK_CARD, bullet="▸", size=15, title_size=16):
    """A card with a title and bullet list."""
    rect(slide, x, y, w, h, bg_col)
    yy = y + 0.12
    if title:
        txt(slide, title, x+0.15, yy, w-0.2, 0.36,
            size=title_size, bold=True, colour=TEAL)
        yy += 0.38
    for item in items:
        txt(slide, f"{bullet}  {item}", x+0.15, yy, w-0.25, 0.4,
            size=size, colour=WHITE)
        yy += 0.38


def slide_number(slide, n, total=25):
    txt(slide, f"{n} / {total}", 12.5, 7.15, 0.8, 0.3,
        size=11, colour=LIGHT_GREY, align=PP_ALIGN.RIGHT)


def footer(slide, text="Opinion Miner  |  Flipkart Review Sentiment Analyser"):
    rect(slide, 0, 7.2, 13.33, 0.3, DARK_CARD)
    txt(slide, text, 0.2, 7.22, 13.0, 0.26,
        size=11, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Gradient-like top strip
rect(s, 0, 0, 13.33, 0.5, DARK_CARD)
rect(s, 0, 0, 13.33, 0.06, TEAL)

# Big title
txt(s, "OPINION MINER", 0.5, 1.3, 12.3, 1.2,
    size=58, bold=True, colour=TEAL, align=PP_ALIGN.CENTER)
txt(s, "Flipkart Review Sentiment Analyser with Fake Review Detection",
    0.5, 2.55, 12.3, 0.7,
    size=22, bold=False, colour=WHITE, align=PP_ALIGN.CENTER)

divider(s, 3.35, TEAL, 0.05)

# Tag chips
chips = [
    ("Python + Flask", MID_BLUE),
    ("Selenium Scraper", MID_BLUE),
    ("VADER NLP", MID_BLUE),
    ("SVM  91.55% Acc", RGBColor(0x00, 0x7A, 0x5E)),
    ("Aspect-Based SA", MID_BLUE),
    ("Docker Deploy", RGBColor(0x20, 0x40, 0x80)),
]
cx = 1.2
for label, col in chips:
    w_ = 1.7
    rect(s, cx, 3.55, w_, 0.42, col)
    txt(s, label, cx, 3.55, w_, 0.42,
        size=13, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    cx += w_ + 0.18

# Subtitle info
txt(s, "A Final Year Project Presentation", 0.5, 4.25, 12.3, 0.4,
    size=17, italic=True, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)
txt(s, "Presented by:  Thagillapally Rakesh", 0.5, 4.8, 12.3, 0.4,
    size=17, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Department of Computer Science & Engineering", 0.5, 5.2, 12.3, 0.4,
    size=15, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Bottom bar
rect(s, 0, 7.1, 13.33, 0.4, RGBColor(0x00, 0xA0, 0x80))
txt(s, "AI • NLP • Machine Learning • Web Development",
    0.5, 7.12, 12.3, 0.35,
    size=13, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
slide_number(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Problem Statement", "Why do we need Opinion Miner?")
footer(s); slide_number(s, 2)

# Large stat boxes
stats = [
    ("87%", "of consumers read online reviews\nbefore buying a product", TEAL),
    ("40%", "of all product reviews on e-commerce\nplatforms are fake or misleading", CORAL),
    ("₹ Crores", "lost by consumers making purchase\ndecisions based on fake reviews", ORANGE),
]
for i, (stat, desc, col) in enumerate(stats):
    rect(s, 0.4 + i*4.3, 1.4, 3.9, 2.2, DARK_CARD)
    rect(s, 0.4 + i*4.3, 1.4, 3.9, 0.07, col)
    txt(s, stat, 0.4 + i*4.3, 1.5, 3.9, 0.9,
        size=42, bold=True, colour=col, align=PP_ALIGN.CENTER)
    txt(s, desc, 0.5 + i*4.3, 2.45, 3.7, 1.1,
        size=13, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Problem description
rect(s, 0.4, 3.85, 12.5, 2.8, DARK_CARD)
txt(s, "Core Challenges:", 0.65, 3.95, 12.0, 0.4,
    size=16, bold=True, colour=TEAL)
problems = [
    "Shoppers cannot distinguish genuine reviews from computer-generated (AI-written) fake ones.",
    "No tool provides both sentiment analysis AND fake review filtering in a single interface.",
    "Existing platforms give only an average star rating — no deep aspect-level insights (Quality, Delivery, Value…).",
    "Manual reading of 100s of reviews is time-consuming and biased.",
]
for i, p in enumerate(problems):
    txt(s, f"  ✗  {p}", 0.65, 4.4 + i*0.52, 12.0, 0.48,
        size=14, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — MOTIVATION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Motivation & Need", "Why this project matters")
footer(s); slide_number(s, 3)

cards = [
    ("🛒 Consumer Trust", TEAL,
     ["Flipkart hosts millions of reviews", "No in-built fake detection", "Users are misled daily"]),
    ("🤖 Rise of AI Fakes", ORANGE,
     ["GPT-generated reviews are indistinguishable", "Sellers abuse reviews for ranking", "Platforms fail to keep up"]),
    ("📊 Data Overload", RGBColor(0x6C, 0x63, 0xFF),
     ["1000+ reviews per product", "No summary of what's good/bad", "Aspect-level gaps ignored"]),
    ("🎯 Research Gap", CORAL,
     ["Existing tools: single sentiment label only", "No combination of fake + sentiment", "No open-source solution for India"]),
]
for i, (title, col, points) in enumerate(cards):
    col_x = 0.4 + (i % 2) * 6.5
    col_y = 1.35 + (i // 2) * 2.85
    rect(s, col_x, col_y, 6.1, 2.6, DARK_CARD)
    rect(s, col_x, col_y, 6.1, 0.06, col)
    txt(s, title, col_x + 0.15, col_y + 0.1, 5.8, 0.45,
        size=17, bold=True, colour=col)
    for j, pt in enumerate(points):
        txt(s, f"▸  {pt}", col_x + 0.15, col_y + 0.6 + j*0.55, 5.8, 0.5,
            size=13, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — OBJECTIVES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Objectives", "What Opinion Miner aims to achieve")
footer(s); slide_number(s, 4)

objectives = [
    ("01", "Scrape Reviews Automatically",
     "Fetch all available reviews from Flipkart product pages using Selenium + BeautifulSoup."),
    ("02", "Detect Fake Reviews (ML)",
     "Filter out AI-generated / fake reviews using a LinearSVC model with 91.55% accuracy."),
    ("03", "Sentiment Classification",
     "Classify each genuine review as Positive, Negative, or Neutral using VADER NLP."),
    ("04", "Aspect-Based Analysis",
     "Break sentiment down into 9 product aspects: Quality, Delivery, Value, Battery, Camera…"),
    ("05", "Visual Analytics",
     "Generate bar charts, score charts and word clouds for quick consumer insights."),
    ("06", "Accessible Interface",
     "Provide a web UI + Chrome Extension so any user can analyse any Flipkart product."),
]
for i, (num, title, desc) in enumerate(objectives):
    col_x = 0.4 + (i % 2) * 6.5
    row_y = 1.3 + (i // 2) * 1.9
    rect(s, col_x, row_y, 6.1, 1.7, DARK_CARD)
    # Number badge
    rect(s, col_x, row_y, 0.55, 1.7, TEAL)
    txt(s, num, col_x, row_y, 0.55, 1.7,
        size=22, bold=True, colour=NAVY, align=PP_ALIGN.CENTER)
    txt(s, title, col_x + 0.65, row_y + 0.1, 5.3, 0.45,
        size=15, bold=True, colour=TEAL)
    txt(s, desc, col_x + 0.65, row_y + 0.58, 5.3, 1.0,
        size=12, colour=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — SCOPE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Scope of the Project", "Inclusions and Exclusions")
footer(s); slide_number(s, 5)

# Included
rect(s, 0.35, 1.3, 5.9, 5.7, DARK_CARD)
rect(s, 0.35, 1.3, 5.9, 0.06, TEAL)
txt(s, "✅  INCLUDED", 0.55, 1.38, 5.5, 0.45,
    size=17, bold=True, colour=TEAL)
included = [
    "Flipkart product review pages (any category)",
    "Manual single-text sentiment analysis",
    "SVM-based fake review detection (ML)",
    "VADER sentiment scoring (Positive/Neutral/Negative)",
    "9-aspect sentiment breakdown",
    "Bar chart, score chart, word cloud visualisations",
    "Chrome Extension for in-browser analysis",
    "Docker-based deployment on Render.com",
    "English language reviews",
]
for i, item in enumerate(included):
    txt(s, f"  ✔  {item}", 0.55, 1.9 + i*0.52, 5.6, 0.5, size=13, colour=WHITE)

# Excluded
rect(s, 6.6, 1.3, 6.35, 5.7, DARK_CARD)
rect(s, 6.6, 1.3, 6.35, 0.06, CORAL)
txt(s, "❌  EXCLUDED", 6.8, 1.38, 6.0, 0.45,
    size=17, bold=True, colour=CORAL)
excluded = [
    "Amazon / other e-commerce platforms",
    "Non-English (multilingual) reviews",
    "Real-time review monitoring / alerts",
    "User login / authentication system",
    "Mobile native app (Android/iOS)",
    "Seller-side analytics dashboard",
    "Social media review analysis",
    "Review summarisation (abstractive NLP)",
]
for i, item in enumerate(excluded):
    txt(s, f"  ✘  {item}", 6.8, 1.9 + i*0.52, 5.9, 0.5, size=13, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — LITERATURE SURVEY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Literature Survey", "Existing systems and research")
footer(s); slide_number(s, 6)

papers = [
    ("Hutto & Gilbert (2014)", "VADER", TEAL,
     "Introduced VADER — a lexicon + rule-based sentiment analyser tuned for social media. Outperforms human raters on micro-blog content."),
    ("Ott et al. (2011)", "Deceptive Opinion Spam", ORANGE,
     "First large-scale study of fake hotel reviews. Used SVM with n-gram features — 90%+ accuracy. Established the CG/OR labelling."),
    ("Mukherjee et al. (2012)", "Yelp Fake Reviewer Groups", CORAL,
     "Detected organised fake reviewer networks on Yelp. Showed behavioural + content features together improve detection."),
    ("Pang & Lee (2008)", "Opinion Mining Survey", RGBColor(0x6C, 0x63, 0xFF),
     "Foundational survey on sentiment analysis. Established aspect-based SA (ABSA) as a key sub-task for fine-grained insights."),
    ("Jindal & Liu (2008)", "Review Spam Detection", RGBColor(0x00, 0xA0, 0xC0),
     "Defined three types of review spam. Used logistic regression on product & reviewer features to identify anomalous reviews."),
]
for i, (auth, topic, col, desc) in enumerate(papers):
    yy = 1.3 + i * 1.18
    rect(s, 0.35, yy, 12.6, 1.08, DARK_CARD)
    rect(s, 0.35, yy, 0.06, 1.08, col)
    txt(s, auth, 0.55, yy + 0.05, 2.5, 0.38, size=13, bold=True, colour=col)
    txt(s, topic, 0.55, yy + 0.48, 2.5, 0.38, size=12, italic=True, colour=LIGHT_GREY)
    txt(s, desc, 3.1, yy + 0.1, 9.6, 0.85, size=12, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — LIMITATIONS OF EXISTING SYSTEMS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Limitations of Existing Systems", "Gaps our project fills")
footer(s); slide_number(s, 7)

rows = [
    ("System / Tool", "What it Does", "Key Limitation", True),
    ("Flipkart Built-in", "Star rating + review display", "No fake detection, no NLP analysis", False),
    ("Amazon Sentiment", "Basic rating histogram", "No aspect breakdown, no authenticity check", False),
    ("TextBlob / VADER standalone", "Sentiment of given text", "No scraper, no fake filter, no visualisation", False),
    ("Fake-Review Detectors (research)", "Classify fake vs genuine", "Offline only, no sentiment pipeline integrated", False),
    ("Opinion Mining Tools (commercial)", "Broad sentiment dashboards", "Expensive, no Flipkart support, no custom training", False),
    ("Opinion Miner  ★ (Our System)", "All-in-one: scrape+filter+SA+aspects+viz", "None — built specifically for Indian e-commerce", False),
]
col_widths = [2.5, 3.5, 4.8]
col_xs     = [0.3, 2.9, 6.5]
for r, row_data in enumerate(rows):
    is_header = row_data[3]
    bg_c = MID_BLUE if is_header else (RGBColor(0x00, 0x50, 0x3A) if r == len(rows)-1 else DARK_CARD)
    yy = 1.28 + r * 0.76
    for c, (cell_text, cw, cx) in enumerate(zip(row_data[:3], col_widths, col_xs)):
        rect(s, cx, yy, cw - 0.04, 0.72, bg_c)
        tc = TEAL if is_header else (TEAL if r == len(rows)-1 and c == 0 else WHITE)
        txt(s, cell_text, cx + 0.1, yy + 0.05, cw - 0.15, 0.62,
            size=13, bold=is_header, colour=tc)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — PROPOSED SYSTEM OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Proposed System — Overview", "A unified end-to-end review intelligence platform")
footer(s); slide_number(s, 8)

# Pipeline flow
steps = [
    ("1\nINPUT", "Flipkart\nProduct URL\nor Text", TEAL),
    ("2\nSCRAPE", "Selenium\nBeautifulSoup\n50-page cap", MID_BLUE),
    ("3\nFILTER", "SVM Fake\nDetector\n91.55% acc", CORAL),
    ("4\nSENTIMENT", "VADER NLP\n+ Negation\nHandling", RGBColor(0x6C, 0x63, 0xFF)),
    ("5\nASPECTS", "9 Product\nAspect\nScores", ORANGE),
    ("6\nOUTPUT", "Charts +\nVerdict +\nDashboard", RGBColor(0x00, 0xA0, 0x60)),
]
for i, (label, desc, col) in enumerate(steps):
    x = 0.35 + i * 2.16
    rect(s, x, 1.5, 1.9, 1.9, col)
    txt(s, label, x, 1.5, 1.9, 0.7,
        size=14, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, x, 2.25, 1.9, 1.1,
        size=12, colour=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, "→", x + 1.93, 2.0, 0.25, 0.8,
            size=28, bold=True, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

# Feature highlights
txt(s, "Key Differentiators:", 0.35, 3.65, 12.6, 0.4,
    size=16, bold=True, colour=TEAL)
diffs = [
    ("🧠  ML-Powered Fake Detection", "LinearSVC trained on 40,432 labelled reviews — not simple rule-based heuristics"),
    ("📐  9-Aspect Sentiment", "Quality · Delivery · Value · Performance · Battery · Camera · Display · Design · Service"),
    ("🌐  Accessible Anywhere", "Web UI + Chrome Extension — works on any device with a browser"),
    ("🐳  Production Ready", "Dockerised, deployed on Render.com — accessible via public URL"),
]
for i, (title, sub) in enumerate(diffs):
    rect(s, 0.35 + (i%2)*6.5, 4.15 + (i//2)*1.45, 6.1, 1.35, DARK_CARD)
    txt(s, title, 0.55 + (i%2)*6.5, 4.2 + (i//2)*1.45, 5.8, 0.45,
        size=14, bold=True, colour=TEAL)
    txt(s, sub, 0.55 + (i%2)*6.5, 4.65 + (i//2)*1.45, 5.8, 0.8,
        size=12, colour=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "System Architecture", "Three-layer design: Presentation → Application → Data")
footer(s); slide_number(s, 9)

# Three layers
layers = [
    ("PRESENTATION LAYER", TEAL, ["Web Browser (HTML/CSS/JS)", "Chrome Extension (Manifest V3)", "User enters URL or text"]),
    ("APPLICATION LAYER",  ORANGE, ["Flask REST API (Python)", "Scraper Module (Selenium+BS4)", "NLP Engine (VADER + SVM)", "Aspect Analyser", "Visualiser (Matplotlib)"]),
    ("DATA / MODEL LAYER", RGBColor(0x6C, 0x63, 0xFF), ["SVM Model (LinearSVC .pkl)", "TF-IDF Vectoriser (.pkl)", "Flipkart Product Reviews (web)", "NLTK Corpora (local)", "Chart images (base64)"]),
]
for i, (lname, col, items) in enumerate(layers):
    x = 0.35 + i * 4.35
    rect(s, x, 1.35, 4.1, 5.55, DARK_CARD)
    rect(s, x, 1.35, 4.1, 0.55, col)
    txt(s, lname, x, 1.35, 4.1, 0.55,
        size=13, bold=True, colour=NAVY, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        rect(s, x + 0.18, 2.0 + j*0.92, 3.74, 0.78, MID_BLUE)
        txt(s, item, x + 0.25, 2.05 + j*0.92, 3.65, 0.7,
            size=13, colour=WHITE, align=PP_ALIGN.CENTER)

# Arrows between layers
for ax in [4.5, 8.85]:
    txt(s, "⇄", ax, 3.3, 0.5, 0.6,
        size=28, bold=True, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Technology Stack", "Tools and frameworks powering Opinion Miner")
footer(s); slide_number(s, 10)

tech = [
    ("🐍 Python 3.11",          "Core language",                   TEAL),
    ("🌶  Flask",               "REST API & web server",           TEAL),
    ("🔵 Selenium",             "Browser automation / scraping",   ORANGE),
    ("🍲 BeautifulSoup 4",      "HTML parsing",                    ORANGE),
    ("📊 VADER Sentiment",       "Lexicon-based NLP",               RGBColor(0x6C,0x63,0xFF)),
    ("🤖 scikit-learn LinearSVC","Fake review detection (ML)",      RGBColor(0x6C,0x63,0xFF)),
    ("🔤 TF-IDF Vectoriser",    "15k-feature text encoding",       RGBColor(0x6C,0x63,0xFF)),
    ("📖 NLTK",                 "Tokenisation, lemmatisation",     RGBColor(0x6C,0x63,0xFF)),
    ("📉 Matplotlib / Seaborn", "Chart generation",                CORAL),
    ("☁  WordCloud",            "Positive/negative word clouds",   CORAL),
    ("🐳 Docker",               "Containerised deployment",        MID_BLUE),
    ("☁  Render.com",           "Cloud hosting (free tier)",       MID_BLUE),
    ("🧩 Chrome Extension MV3", "Browser-level integration",       ORANGE),
    ("🔧 webdriver-manager",    "Auto ChromeDriver management",    MID_BLUE),
]
for i, (tool, role, col) in enumerate(tech):
    cx = 0.35 + (i % 2) * 6.5
    cy = 1.32 + (i // 2) * 0.75
    rect(s, cx, cy, 6.1, 0.68, DARK_CARD)
    rect(s, cx, cy, 0.06, 0.68, col)
    txt(s, tool, cx + 0.18, cy + 0.04, 2.5, 0.6, size=13, bold=True, colour=col)
    txt(s, role, cx + 2.7, cy + 0.04, 3.3, 0.6, size=12, colour=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — MODULE BREAKDOWN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Module Breakdown", "Six core components")
footer(s); slide_number(s, 11)

modules = [
    ("scraper/\nflipkart_scraper.py", TEAL, [
        "Selenium headless Chrome", "Multi-page scrolling (50 pages)", "Junk review filter",
        "RNW layout parser", "Snap-path bypass for Chrome"]),
    ("nlp/\nsentiment.py", RGBColor(0x6C,0x63,0xFF), [
        "VADER polarity scores", "Custom negation handling",
        "Compound → Pos/Neu/Neg label", "Returns scores dict"]),
    ("nlp/\nfake_detector.py", CORAL, [
        "LinearSVC model (91.55%)", "TF-IDF 15k features + bigrams",
        "Lemmatisation + stop-words", "Fallback rule-based detector"]),
    ("nlp/\naspect_sentiment.py", ORANGE, [
        "9 aspects via keyword map", "Sentence-level tokenisation",
        "VADER per sentence", "avg_compound + verdict"]),
    ("utils/\naggregator.py", RGBColor(0x00,0xA0,0xC0), [
        "Counts Pos/Neu/Neg/Fake", "Average scores computation",
        "Overall product verdict", "Action recommendation"]),
    ("utils/\nvisualizer.py", RGBColor(0x80,0x60,0xFF), [
        "Sentiment bar chart (Matplotlib)", "Score bar chart",
        "Word cloud — Positive", "Word cloud — Negative"]),
]
for i, (name, col, points) in enumerate(modules):
    cx = 0.3 + (i % 3) * 4.38
    cy = 1.35 + (i // 3) * 3.0
    rect(s, cx, cy, 4.1, 2.75, DARK_CARD)
    rect(s, cx, cy, 4.1, 0.06, col)
    txt(s, name, cx + 0.1, cy + 0.1, 3.9, 0.65,
        size=13, bold=True, colour=col)
    for j, pt in enumerate(points):
        txt(s, f"▸  {pt}", cx + 0.1, cy + 0.82 + j*0.42, 3.9, 0.4,
            size=11, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — DATA FLOW DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Data Flow & Workflow", "End-to-end journey of a review analysis request")
footer(s); slide_number(s, 12)

# DFD-style flow
flow = [
    ("User enters\nFlipkart URL",      0.3,  1.5,  TEAL),
    ("Flask /analyze-url\nreceives POST", 3.1, 1.5, MID_BLUE),
    ("Selenium scrapes\npages 1–50",    5.9,  1.5,  ORANGE),
    ("Reviews extracted\nvia BS4 parser",8.7, 1.5,  ORANGE),
    ("SVM Fake Detector\nfilters reviews",1.7, 3.6, CORAL),
    ("VADER analyses\ngenuine reviews", 4.5,  3.6,  RGBColor(0x6C,0x63,0xFF)),
    ("Aspect Analyser\n9 dimensions",   7.3,  3.6,  RGBColor(0x00,0xA0,0x60)),
    ("Visualiser creates\ncharts (b64)", 10.1, 3.6, RGBColor(0x80,0x60,0xFF)),
    ("JSON response\nto frontend",      4.5,  5.7,  TEAL),
    ("Charts + Verdict\ndisplayed to user", 8.0, 5.7, RGBColor(0x00,0xD4,0x80)),
]
for (label, x, y, col) in flow:
    rect(s, x, y, 2.5, 1.1, col)
    txt(s, label, x, y, 2.5, 1.1,
        size=12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

# Connecting arrows (text-based)
arrows = [
    (2.82, 1.93, "→"), (5.62, 1.93, "→"), (8.42, 1.93, "→"),
    (1.7, 2.63, "↓"), (4.5, 2.63, "↓"), (7.3, 2.63, "↓"), (10.1, 2.63, "↓"),
    (4.22, 4.12, "→"), (7.02, 4.12, "→"),
    (4.5, 4.73, "↓"), (8.0, 4.73, "↓"),
    (6.52, 6.12, "→"),
]
for (ax, ay, arrow) in arrows:
    txt(s, arrow, ax, ay, 0.4, 0.45,
        size=20, bold=True, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Methodology & Approach", "Step-by-step implementation strategy")
footer(s); slide_number(s, 13)

steps2 = [
    ("Phase 1", "Data Collection",
     "Selenium WebDriver navigates Flipkart review pages in headless mode. "
     "BeautifulSoup parses the React Native Web (RNW) DOM. "
     "Up to 50 pages scraped per product URL.", TEAL),
    ("Phase 2", "Preprocessing",
     "Reviews cleaned: lowercasing, stop-word removal, lemmatisation (NLTK). "
     "Junk reviews filtered by keyword blacklists (address, contact info). "
     "Duplicate exact-match detection.", ORANGE),
    ("Phase 3", "Fake Review Detection",
     "TF-IDF (15,000 unigram + bigram features) vectorises each review. "
     "Pre-trained LinearSVC predicts Genuine (OR) vs Fake (CG). "
     "Only Genuine reviews passed to sentiment pipeline.", CORAL),
    ("Phase 4", "Sentiment Analysis",
     "VADER polarity_scores() applied to each genuine review. "
     "Custom negation handler adjusts compound score. "
     "Compound > 0.2 → Positive; < -0.2 → Negative; else Neutral.", RGBColor(0x6C,0x63,0xFF)),
    ("Phase 5", "Aspect-Based Analysis",
     "Each review split into sentences (NLTK sent_tokenize). "
     "Sentences matched to 9 aspect keyword lists. "
     "VADER applied per sentence → per-aspect avg_compound + verdict.", RGBColor(0x00,0xA0,0x60)),
    ("Phase 6", "Aggregation & Visualisation",
     "Counts, averages, and overall verdict computed. "
     "Matplotlib generates bar chart, score chart; WordCloud generates word clouds. "
     "All charts encoded as base64 PNG and sent in JSON response.", MID_BLUE),
]
for i, (phase, title, body, col) in enumerate(steps2):
    cx = 0.3 + (i % 2) * 6.52
    cy = 1.3 + (i // 2) * 2.02
    rect(s, cx, cy, 6.15, 1.88, DARK_CARD)
    rect(s, cx, cy, 0.55, 1.88, col)
    txt(s, phase, cx, cy, 0.55, 1.88,
        size=12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, cx + 0.65, cy + 0.08, 5.35, 0.42,
        size=14, bold=True, colour=col)
    txt(s, body, cx + 0.65, cy + 0.55, 5.35, 1.25,
        size=11, colour=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — ALGORITHMS / MODELS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Algorithms & Models Used", "The intelligence behind Opinion Miner")
footer(s); slide_number(s, 14)

# VADER box
rect(s, 0.3, 1.35, 6.1, 5.7, DARK_CARD)
rect(s, 0.3, 1.35, 6.1, 0.06, TEAL)
txt(s, "VADER — Sentiment Analysis", 0.5, 1.42, 5.8, 0.5,
    size=16, bold=True, colour=TEAL)
vader_points = [
    "Type: Lexicon + Rule-Based NLP",
    "Valence Aware Dictionary for Sentiment Reasoning",
    "Returns: pos, neg, neu, compound scores",
    "Compound range: −1.0 (very negative) to +1.0 (very positive)",
    "Thresholds used:",
    "   compound > 0.2  →  POSITIVE",
    "   compound < −0.2 →  NEGATIVE",
    "   otherwise       →  NEUTRAL",
    "Custom Negation Handling:",
    '   "not good" / "isn\'t working" → compound adjusted',
    "Why VADER? Fast, no training needed,",
    "           works on short review text",
]
for i, pt in enumerate(vader_points):
    bold = "→" in pt or ":" in pt and i < 5
    txt(s, pt, 0.5, 2.0 + i*0.41, 5.8, 0.4,
        size=12, colour=WHITE if "Why" not in pt else LIGHT_GREY,
        bold=bold)

# SVM box
rect(s, 6.75, 1.35, 6.22, 5.7, DARK_CARD)
rect(s, 6.75, 1.35, 6.22, 0.06, CORAL)
txt(s, "LinearSVC — Fake Review Detector", 6.95, 1.42, 5.9, 0.5,
    size=16, bold=True, colour=CORAL)
svm_points = [
    "Type: Support Vector Machine (Linear kernel)",
    "Training data: 40,432 reviews (balanced)",
    "   CG = 20,216 (Computer Generated / Fake)",
    "   OR = 20,216 (Original Reviews / Genuine)",
    "Features: TF-IDF",
    "   Max features : 15,000",
    "   N-gram range : (1,2) unigrams + bigrams",
    "   Sublinear TF : True",
    "Train/Test split: 80% / 20%",
    "Performance Metrics:",
    "   Accuracy  : 91.55%",
    "   Precision : 91.63%",
    "   Recall    : 91.47%",
    "   F1-Score  : 91.55%",
]
for i, pt in enumerate(svm_points):
    col = YELLOW if any(x in pt for x in ["91", "Accuracy", "Precision", "Recall", "F1"]) else WHITE
    txt(s, pt, 6.95, 2.0 + i*0.41, 5.9, 0.4,
        size=12, colour=col,
        bold=any(x in pt for x in ["91.55", "91.63", "91.47"]))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — DATASET
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Dataset Description", "Training data for the fake review detector")
footer(s); slide_number(s, 15)

# Dataset stats
stats2 = [
    ("40,432", "Total Reviews", TEAL),
    ("20,216", "Genuine (OR)", RGBColor(0x00,0xC8,0x70)),
    ("20,216", "Fake (CG)", CORAL),
    ("80% / 20%", "Train / Test Split", ORANGE),
]
for i, (num, label, col) in enumerate(stats2):
    rect(s, 0.3 + i*3.23, 1.35, 3.05, 1.5, DARK_CARD)
    rect(s, 0.3 + i*3.23, 1.35, 3.05, 0.06, col)
    txt(s, num, 0.3 + i*3.23, 1.45, 3.05, 0.75,
        size=32, bold=True, colour=col, align=PP_ALIGN.CENTER)
    txt(s, label, 0.3 + i*3.23, 2.2, 3.05, 0.5,
        size=14, colour=WHITE, align=PP_ALIGN.CENTER)

# Details
rect(s, 0.3, 3.05, 12.7, 3.85, DARK_CARD)
details = [
    ("Source", "Ott et al. (Cornell) Deceptive Opinion Spam dataset — widely used benchmark in NLP research"),
    ("Labels", "CG (Computer Generated) → Fake = 1  |  OR (Original Reviews) → Genuine = 0"),
    ("Preprocessing", "Lowercasing → URL removal → non-alpha removal → Word tokenise → Lemmatise → Stop-words filter"),
    ("Balance", "Perfectly balanced dataset (50/50) eliminates class-imbalance bias — no SMOTE needed"),
    ("Why SVM?", "LinearSVC is proven best-in-class for high-dimensional sparse TF-IDF matrices (faster than neural nets, comparable accuracy)"),
    ("Validation", "Stratified 80/20 split ensures equal class ratio in both train and test sets"),
]
for i, (key, val) in enumerate(details):
    txt(s, f"{key}:", 0.55, 3.18 + i*0.58, 1.8, 0.5,
        size=13, bold=True, colour=TEAL)
    txt(s, val, 2.4, 3.18 + i*0.58, 10.35, 0.5,
        size=12, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — KEY FEATURES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Key Features Implemented", "What makes Opinion Miner stand out")
footer(s); slide_number(s, 16)

features = [
    ("🔍 Flipkart Scraper", TEAL, [
        "Headless Chrome via Selenium", "Parses Flipkart's React layout", "Handles dynamic JS content", "Auto ChromeDriver via webdriver-manager"]),
    ("🤖 SVM Fake Detector", CORAL, [
        "91.55% accuracy", "TF-IDF + bigrams", "Pre-trained — no runtime training", "Instant classification per review"]),
    ("📊 VADER Sentiment", RGBColor(0x6C,0x63,0xFF), [
        "Pos/Neu/Neg + compound score", "Custom negation handling", "Avg scores shown for scraped reviews", "Emoji + audio feedback"]),
    ("📐 Aspect Sentiment", ORANGE, [
        "9 product aspects analysed", "Sentence-level granularity", "Colour-coded verdict cards", "Sample sentences per aspect"]),
    ("📈 Visualisations", RGBColor(0x00,0xA0,0xC0), [
        "Sentiment distribution bar chart", "Pos/Neu/Neg score chart", "Positive word cloud", "Negative word cloud"]),
    ("🧩 Chrome Extension", RGBColor(0x80,0x60,0xFF), [
        "Manifest V3 extension", "Works on any Flipkart page", "Shows results in popup", "Calls local Flask API"]),
]
for i, (title, col, pts) in enumerate(features):
    cx = 0.3 + (i % 3) * 4.35
    cy = 1.32 + (i // 3) * 3.05
    rect(s, cx, cy, 4.1, 2.8, DARK_CARD)
    rect(s, cx, cy, 4.1, 0.06, col)
    txt(s, title, cx + 0.1, cy + 0.1, 3.9, 0.5, size=14, bold=True, colour=col)
    for j, pt in enumerate(pts):
        txt(s, f"✔  {pt}", cx + 0.1, cy + 0.68 + j*0.5, 3.9, 0.48, size=12, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — SCREENSHOTS / APPLICATION UI
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Application Screenshots", "Opinion Miner in action")
footer(s); slide_number(s, 17)

screens = [
    ("🏠 Home Page", "Clean dark-themed UI with two modes:\n• Manual text input\n• Flipkart URL analysis", TEAL),
    ("📝 Text Analysis", "Enter any review text → instant\nsentiment + emoji + audio feedback\nCompound score breakdown shown", RGBColor(0x6C,0x63,0xFF)),
    ("🔗 URL Analysis", "Paste Flipkart URL → scrape → filter\nfakes → analyse → full dashboard\nwith all metrics in one view", ORANGE),
    ("📊 Charts Panel", "Sentiment bar chart, score chart,\nPositive & Negative word clouds\ngenerated automatically", CORAL),
    ("📐 Aspect Cards", "9 colour-coded aspect cards showing\nPositive/Neutral/Negative counts +\navg compound score per aspect", RGBColor(0x00,0xA0,0x60)),
    ("🧩 Chrome Extension", "Browser popup: analyse any Flipkart\npage without leaving the tab.\nShows verdict + aspect breakdown", MID_BLUE),
]
for i, (title, desc, col) in enumerate(screens):
    cx = 0.28 + (i % 3) * 4.35
    cy = 1.32 + (i // 3) * 3.05
    rect(s, cx, cy, 4.1, 2.8, DARK_CARD)
    rect(s, cx, cy, 4.1, 0.06, col)
    # Placeholder screen area
    rect(s, cx + 0.15, cy + 0.55, 3.8, 1.5, MID_BLUE)
    txt(s, "[ Screenshot ]", cx + 0.15, cy + 0.55, 3.8, 1.5,
        size=14, colour=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)
    txt(s, title, cx + 0.1, cy + 0.1, 3.9, 0.42, size=13, bold=True, colour=col)
    txt(s, desc, cx + 0.1, cy + 2.12, 3.9, 0.65, size=11, colour=LIGHT_GREY)

txt(s, "⚠  Replace placeholder boxes with actual screenshots before presentation",
    0.3, 7.0, 12.7, 0.35,
    size=11, colour=ORANGE, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — PERFORMANCE METRICS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Performance Metrics & Evaluation", "How well does Opinion Miner perform?")
footer(s); slide_number(s, 18)

# Big metric boxes
metrics = [
    ("91.55%", "Accuracy", TEAL),
    ("91.63%", "Precision", RGBColor(0x00,0xC8,0x70)),
    ("91.47%", "Recall",    ORANGE),
    ("91.55%", "F1-Score",  CORAL),
]
for i, (val, label, col) in enumerate(metrics):
    rect(s, 0.3 + i*3.23, 1.35, 3.05, 1.8, DARK_CARD)
    rect(s, 0.3 + i*3.23, 1.35, 3.05, 0.08, col)
    txt(s, val, 0.3 + i*3.23, 1.5, 3.05, 0.95,
        size=38, bold=True, colour=col, align=PP_ALIGN.CENTER)
    txt(s, label, 0.3 + i*3.23, 2.46, 3.05, 0.5,
        size=15, colour=WHITE, align=PP_ALIGN.CENTER)

# Classification breakdown
rect(s, 0.3, 3.35, 12.7, 1.12, DARK_CARD)
txt(s, "Classification Report — Test Set (8,086 reviews)", 0.5, 3.42, 12.3, 0.4,
    size=14, bold=True, colour=TEAL)
headers = ["Class", "Precision", "Recall", "F1", "Support"]
widths  = [2.5, 2.2, 2.2, 2.2, 2.2]
xs      = [0.5, 3.1, 5.35, 7.6, 9.85]
for h, x, w in zip(headers, xs, widths):
    txt(s, h, x, 3.85, w, 0.38, size=12, bold=True, colour=LIGHT_GREY)

rows_data = [("Genuine (OR)", "91.47%", "91.63%", "91.55%", "4,043"),
             ("Fake (CG)",    "91.63%", "91.47%", "91.55%", "4,043")]
for r, row in enumerate(rows_data):
    col_ = TEAL if r == 0 else CORAL
    for val, x, w in zip(row, xs, widths):
        txt(s, val, x, 4.25 + r*0.38, w, 0.36, size=12,
            colour=col_ if val in ["Genuine (OR)", "Fake (CG)"] else WHITE)

# Scraper speed
rect(s, 0.3, 5.0, 12.7, 2.3, DARK_CARD)
txt(s, "Scraper & System Performance", 0.5, 5.08, 12.3, 0.42,
    size=14, bold=True, colour=TEAL)
perf = [
    ("Pages scraped per run", "Up to 50 (configurable)", ORANGE),
    ("Avg reviews per page", "~10 reviews", ORANGE),
    ("Typical scrape time (50 pages)", "60–90 seconds", ORANGE),
    ("Fake detection speed", "<1 ms per review (LinearSVC)", RGBColor(0x00,0xC8,0x70)),
    ("VADER sentiment speed", "<1 ms per review", RGBColor(0x00,0xC8,0x70)),
    ("Chart generation time", "<2 seconds (all 4 charts)", RGBColor(0x00,0xC8,0x70)),
]
for i, (metric, val, col) in enumerate(perf):
    cx = 0.5 + (i % 2) * 6.4
    cy = 5.55 + (i // 2) * 0.55
    txt(s, f"▸  {metric}:", cx, cy, 3.8, 0.48, size=12, colour=LIGHT_GREY)
    txt(s, val, cx + 3.85, cy, 2.4, 0.48, size=12, bold=True, colour=col)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — COMPARISON WITH EXISTING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Comparison with Existing Systems", "Opinion Miner vs alternatives")
footer(s); slide_number(s, 19)

cols_data = ["Feature", "Flipkart\nBuilt-in", "TextBlob /\nVADER Standalone", "Research\nPapers", "Opinion\nMiner ★"]
col_xs2   = [0.28, 2.8, 5.05, 7.35, 9.7]
col_ws2   = [2.45, 2.2, 2.2, 2.25, 3.35]
features2 = [
    ("Flipkart Review Scraping",   ["✓", "✗", "✗", "✓"]),
    ("Fake Review Detection",      ["✗", "✗", "✓", "✓"]),
    ("VADER Sentiment Analysis",   ["✗", "✓", "✗", "✓"]),
    ("Aspect-Level Sentiment",     ["✗", "✗", "Partial", "✓"]),
    ("Visual Charts / Word Cloud", ["✗", "✗", "✗", "✓"]),
    ("Chrome Extension",           ["✗", "✗", "✗", "✓"]),
    ("Docker / Cloud Deployed",    ["✓", "✗", "✗", "✓"]),
    ("Free & Open Source",         ["✗", "✓", "Partial", "✓"]),
    ("No Setup Required (Web UI)", ["✓", "✗", "✗", "✓"]),
]

# Header row
for txt_val, x, w in zip(cols_data, col_xs2, col_ws2):
    rect(s, x, 1.3, w - 0.04, 0.7, MID_BLUE)
    txt(s, txt_val, x, 1.3, w - 0.04, 0.7,
        size=12, bold=True, colour=TEAL, align=PP_ALIGN.CENTER)

for r, (feat, vals) in enumerate(features2):
    yy = 2.05 + r * 0.56
    bg_c = DARK_CARD if r % 2 == 0 else RGBColor(0x14, 0x28, 0x3E)
    rect(s, col_xs2[0], yy, col_ws2[0] - 0.04, 0.52, bg_c)
    txt(s, feat, col_xs2[0] + 0.1, yy + 0.06, col_ws2[0] - 0.14, 0.42, size=12, colour=WHITE)
    for c, (val, x, w) in enumerate(zip(vals, col_xs2[1:], col_ws2[1:])):
        last = (c == 3)
        bg_cc = RGBColor(0x00, 0x40, 0x30) if last else bg_c
        rect(s, x, yy, w - 0.04, 0.52, bg_cc)
        vc = RGBColor(0x00,0xC8,0x70) if val == "✓" else (CORAL if val == "✗" else YELLOW)
        txt(s, val, x, yy + 0.06, w - 0.04, 0.42,
            size=13, bold=last, colour=vc, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — ADVANTAGES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Advantages of Opinion Miner", "Why this system is better")
footer(s); slide_number(s, 20)

advantages = [
    ("🎯 High Accuracy", TEAL,
     "91.55% F1-score on fake review detection using LinearSVC — outperforms simple rule-based methods significantly."),
    ("⚡ Speed", ORANGE,
     "Fake detection and sentiment analysis complete in <1 ms per review. Full 50-page analysis in under 2 minutes."),
    ("🔬 Granular Insights", RGBColor(0x6C,0x63,0xFF),
     "9-aspect breakdown reveals exactly which product dimensions (Quality, Battery, Camera…) are praised or criticised."),
    ("🌐 Zero Setup", RGBColor(0x00,0xA0,0x60),
     "Users just paste a URL — no API keys, no registration, no installation. Works entirely in the browser."),
    ("🔓 Open Architecture", CORAL,
     "Modular Flask + Python design: easy to extend with new platforms, new NLP models, or new aspects."),
    ("🐳 Production Grade", MID_BLUE,
     "Dockerised deployment on Render.com. Environment variables for ChromeDriver. Gunicorn WSGI server."),
    ("📦 Integrated Pipeline", YELLOW,
     "Scrape → Filter → Analyse → Visualise — all in one request. No need to chain separate tools."),
    ("🧩 Multi-Interface", RGBColor(0x80,0x60,0xFF),
     "Web UI + Chrome Extension: users can analyse from the app or directly from the Flipkart product page."),
]
for i, (title, col, desc) in enumerate(advantages):
    cx = 0.3 + (i % 2) * 6.52
    cy = 1.32 + (i // 2) * 1.5
    rect(s, cx, cy, 6.15, 1.38, DARK_CARD)
    rect(s, cx, cy, 0.06, 1.38, col)
    txt(s, title, cx + 0.18, cy + 0.08, 5.8, 0.42, size=14, bold=True, colour=col)
    txt(s, desc, cx + 0.18, cy + 0.55, 5.8, 0.75, size=12, colour=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — CHALLENGES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Challenges Faced", "Real problems encountered and how we solved them")
footer(s); slide_number(s, 21)

challenges = [
    ("Chrome Snap Conflict",
     "Snap-packaged Chromium on Ubuntu has sandbox restrictions incompatible with headless mode, causing session crash.",
     "Detect and skip /snap/ paths. Prioritise webdriver-manager to auto-download matching ChromeDriver binary."),
    ("Flipkart RNW Layout",
     "Flipkart renders reviews using React Native Web — dynamic JS with no static HTML. Standard BS4 parsing fails.",
     "Selenium waits for JS rendering; custom span.css-1jxf684 selector targets the RNW text nodes directly."),
    ("Fake Review Dataset",
     "No Flipkart-specific labelled fake review dataset available publicly for Indian e-commerce.",
     "Used Cornell Deceptive Opinion Spam dataset (40,432 reviews) — proven benchmark; model generalises well."),
    ("Dedup vs Short Reviews",
     "Aggressive substring deduplication was dropping legitimate short reviews (e.g., 'key ring works perfectly').",
     "Switched to exact-match deduplication only. Removed overly broad junk keywords (ring, block, floor)."),
    ("Chrome 138 Flag Crashes",
     "--js-flags=--max-old-space-size=256 and experimental options crash Chrome 138 with SessionNotCreated error.",
     "Removed incompatible flags. Used only stable flags: --headless=new, --no-sandbox, --disable-gpu."),
    ("Page Load Time",
     "Scraping 100+ pages takes 5+ minutes, making the user wait too long for results.",
     "Capped at 50 pages (~400–500 reviews) as a practical balance between coverage and speed."),
]
for i, (ch, prob, sol) in enumerate(challenges):
    cy = 1.32 + i * 1.02
    rect(s, 0.28, cy, 12.75, 0.95, DARK_CARD)
    rect(s, 0.28, cy, 0.06, 0.95, CORAL)
    txt(s, f"⚠  {ch}", 0.45, cy + 0.04, 3.0, 0.38, size=13, bold=True, colour=CORAL)
    txt(s, f"Problem: {prob}", 3.55, cy + 0.04, 5.5, 0.42, size=11, colour=LIGHT_GREY)
    txt(s, f"✔ {sol}", 3.55, cy + 0.5, 9.3, 0.42, size=11, colour=RGBColor(0x00,0xC8,0x70))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 22 — LIMITATIONS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Limitations", "Known constraints of the current system")
footer(s); slide_number(s, 22)

limitations = [
    ("Flipkart Only", CORAL,
     "Only Flipkart is supported. Amazon, Meesho, Myntra URLs are not scraped. Extending requires site-specific parsers."),
    ("English Only", ORANGE,
     "VADER is English-tuned. Reviews in Telugu, Hindi, Tamil etc. are not correctly analysed (multilingual support removed for device constraints)."),
    ("50-Page Cap", YELLOW,
     "Maximum 50 pages scraped per run (~400–500 reviews). Products with 5000+ reviews will not be fully covered."),
    ("No Real-Time Monitoring", LIGHT_GREY,
     "System analyses reviews on demand only — no scheduled monitoring or alerts for rating changes over time."),
    ("Static Keyword Aspects", RGBColor(0x6C,0x63,0xFF),
     "Aspect detection uses a fixed keyword list. Niche products with unusual terminology may not map to any aspect."),
    ("Fake Model Domain Shift", RGBColor(0x00,0xA0,0xC0),
     "SVM trained on hotel/product reviews from Cornell dataset. May have lower accuracy on very domain-specific Flipkart reviews."),
    ("Scraper Fragility", MID_BLUE,
     "Flipkart UI changes (DOM restructuring) will break the scraper. Requires manual selector updates on each redesign."),
    ("No User Accounts", LIGHT_GREY,
     "No history, saved analyses, or personalised recommendations. Each session is stateless."),
]
for i, (title, col, desc) in enumerate(limitations):
    cx = 0.3 + (i % 2) * 6.52
    cy = 1.32 + (i // 2) * 1.5
    rect(s, cx, cy, 6.15, 1.38, DARK_CARD)
    rect(s, cx, cy, 0.06, 1.38, col)
    txt(s, f"⚠  {title}", cx + 0.18, cy + 0.08, 5.8, 0.42, size=13, bold=True, colour=col)
    txt(s, desc, cx + 0.18, cy + 0.55, 5.8, 0.75, size=12, colour=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 23 — FUTURE SCOPE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Future Scope & Improvements", "Where Opinion Miner can go next")
footer(s); slide_number(s, 23)

future = [
    ("🌐 Multi-Platform Support", TEAL,
     "Add scrapers for Amazon India, Meesho, Myntra. Create a unified product comparison dashboard."),
    ("🗣  Multilingual NLP", RGBColor(0x6C,0x63,0xFF),
     "Integrate XLM-RoBERTa or MuRIL for Hindi/Telugu/Tamil reviews using a GPU-enabled server."),
    ("📱 Mobile Application", ORANGE,
     "React Native app that scans product QR codes and shows instant sentiment analysis."),
    ("🔔 Real-Time Alerts", CORAL,
     "Scheduled monitoring: alert users when a product's sentiment drops below a threshold."),
    ("🧠 Transformer Models", RGBColor(0x00,0xA0,0xC0),
     "Replace VADER with fine-tuned BERT/RoBERTa for higher sentiment accuracy on product reviews."),
    ("📊 Seller Dashboard", YELLOW,
     "Analytics portal for sellers: track how genuine buyers perceive each product aspect over time."),
    ("🤖 Review Summarisation", RGBColor(0x80,0x60,0xFF),
     "Abstractive summarisation using GPT-4 or T5 — generate a 3-line summary of 1000 reviews."),
    ("🏪 Marketplace API", RGBColor(0x00,0xC8,0x70),
     "Official Flipkart Affiliate API integration to replace Selenium — more stable, no scraper fragility."),
]
for i, (title, col, desc) in enumerate(future):
    cx = 0.3 + (i % 2) * 6.52
    cy = 1.32 + (i // 2) * 1.5
    rect(s, cx, cy, 6.15, 1.38, DARK_CARD)
    rect(s, cx, cy, 0.06, 1.38, col)
    txt(s, title, cx + 0.18, cy + 0.08, 5.8, 0.42, size=14, bold=True, colour=col)
    txt(s, desc, cx + 0.18, cy + 0.55, 5.8, 0.75, size=12, colour=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 24 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Conclusion", "Summary of what was built and achieved")
footer(s); slide_number(s, 24)

rect(s, 0.3, 1.35, 12.7, 2.0, DARK_CARD)
rect(s, 0.3, 1.35, 12.7, 0.07, TEAL)
txt(s, "Opinion Miner is a complete, production-ready review intelligence platform that:\n"
    "addresses the critical problem of fake reviews and information overload on Flipkart, "
    "using proven Machine Learning and NLP techniques to deliver actionable purchase insights.",
    0.55, 1.45, 12.2, 1.8, size=15, colour=WHITE)

# Achievements
txt(s, "Key Achievements:", 0.3, 3.55, 12.7, 0.42, size=16, bold=True, colour=TEAL)
achievements = [
    ("✅ Built a real Flipkart scraper", "capable of fetching hundreds of reviews from any product URL"),
    ("✅ Trained LinearSVC model",       "achieving 91.55% accuracy on 40,432 fake review detection samples"),
    ("✅ Integrated VADER NLP",          "with custom negation handling for accurate English review sentiment"),
    ("✅ Implemented 9-aspect ABSA",     "giving granular insight: Quality, Delivery, Battery, Camera and more"),
    ("✅ Created Chrome Extension",      "allowing in-browser Flipkart analysis without visiting the web app"),
    ("✅ Deployed on Render.com",        "using Docker — publicly accessible with zero server management"),
]
for i, (bold_part, normal_part) in enumerate(achievements):
    cy = 4.05 + i * 0.52
    txt(s, bold_part, 0.5, cy, 3.6, 0.48, size=13, bold=True, colour=RGBColor(0x00,0xC8,0x70))
    txt(s, normal_part, 4.15, cy, 8.6, 0.48, size=13, colour=WHITE)

rect(s, 0.3, 7.0, 12.7, 0.38, RGBColor(0x00, 0xA0, 0x80))
txt(s, "\"A single tool that tells you: Is it real? How do people feel? What's good or bad? — for any Flipkart product.\"",
    0.5, 7.02, 12.3, 0.35, size=12, italic=True, colour=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 25 — REFERENCES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "References", "Academic papers and technical resources")
footer(s); slide_number(s, 25)

references = [
    "[1]  Hutto, C.J. & Gilbert, E. (2014). VADER: A Parsimonious Rule-based Model for Sentiment Analysis of Social Media Text. ICWSM.",
    "[2]  Ott, M., Choi, Y., Cardie, C., & Hancock, J.T. (2011). Finding Deceptive Opinion Spam by Any Stretch of the Imagination. ACL.",
    "[3]  Mukherjee, A., Liu, B., & Glance, N. (2012). Spotting Fake Reviewer Groups in Consumer Reviews. WWW Conference.",
    "[4]  Pang, B., & Lee, L. (2008). Opinion Mining and Sentiment Analysis. Foundations and Trends in Information Retrieval.",
    "[5]  Jindal, N., & Liu, B. (2008). Opinion Spam and Analysis. ACM WSDM Conference.",
    "[6]  scikit-learn developers. LinearSVC Documentation. https://scikit-learn.org/stable/modules/svm.html",
    "[7]  NLTK Project. Natural Language Toolkit. https://www.nltk.org",
    "[8]  Selenium Project. WebDriver Documentation. https://www.selenium.dev/documentation",
    "[9]  Mitchell, R. (2018). Web Scraping with Python. O'Reilly Media.",
    "[10] Docker Inc. Docker Documentation. https://docs.docker.com",
    "[11] Render.com. Render Cloud Hosting Platform. https://render.com",
    "[12] Google. Chrome for Developers — Manifest V3. https://developer.chrome.com/docs/extensions/mv3",
]
rect(s, 0.3, 1.32, 12.7, 5.6, DARK_CARD)
for i, ref in enumerate(references):
    col = TEAL if i % 2 == 0 else LIGHT_GREY
    txt(s, ref, 0.5, 1.45 + i * 0.44, 12.3, 0.42, size=12, colour=col)

rect(s, 0.3, 7.0, 12.7, 0.38, DARK_CARD)
txt(s, "Thank You for your attention  —  Questions are welcome! 🙏",
    0.5, 7.02, 12.3, 0.36,
    size=14, bold=True, colour=TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = "/home/rakesh/Desktop/RockySIrProject/Opinion_Miner_Presentation.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
